import os
import pdfplumber
import PyPDF2
from docx import Document
from pptx import Presentation
import streamlit as st


# i wrote this to handle all the different file types we might get
def extract_text_from_pdf(file):
    text = ""
    try:
        with pdfplumber.open(file) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
    except Exception as e:
        st.warning(f"pdfplumber failed, trying PyPDF2... ({e})")
        try:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                text += page.extract_text() + "\n"
        except Exception as e2:
            st.error(f"Could not read PDF: {e2}")
    return text


def extract_text_from_docx(file):
    text = ""
    try:
        doc = Document(file)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        st.error(f"Could not read DOCX: {e}")
    return text


def extract_text_from_pptx(file):
    text = ""
    try:
        prs = Presentation(file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        st.error(f"Could not read PPTX: {e}")
    return text


def detect_content_type(filename, text):
    """
    try to figure out what kind of document this is
    based on filename and content keywords
    """
    filename_lower = filename.lower()
    text_lower = text.lower()[:500]  # only check start of doc

    if any(word in filename_lower for word in ["question", "qb", "exam", "paper", "pyq"]):
        return "question_paper"
    elif any(word in filename_lower for word in ["lab", "practical", "experiment"]):
        return "lab_manual"
    elif any(word in filename_lower for word in ["note", "lecture", "slide"]):
        return "notes"
    elif any(word in text_lower for word in ["chapter", "unit", "introduction", "definition"]):
        return "textbook"
    else:
        return "notes"  # default


def process_uploaded_file(uploaded_file):
    """main function - takes a streamlit uploaded file and returns text + metadata"""
    filename = uploaded_file.name
    extension = filename.split(".")[-1].lower()

    text = ""
    if extension == "pdf":
        text = extract_text_from_pdf(uploaded_file)
    elif extension == "docx":
        text = extract_text_from_docx(uploaded_file)
    elif extension == "pptx":
        text = extract_text_from_pptx(uploaded_file)
    else:
        st.error(f"Unsupported file type: {extension}")
        return None

    if not text.strip():
        st.warning(f"No text found in {filename}. It might be a scanned/image PDF.")
        return None

    content_type = detect_content_type(filename, text)

    return {
        "filename": filename,
        "content_type": content_type,
        "text": text,
        "char_count": len(text),
        "extension": extension
    }